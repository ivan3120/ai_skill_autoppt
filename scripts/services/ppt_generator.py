# -*- coding: utf-8 -*-
"""
PPT生成服务 V2
支持丰富的PPT故事线：整体评估概览、每个MT的详细结果、风险、优化建议
"""

import os
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.table import Table
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


class PPTGeneratorError(Exception):
    """PPT生成相关异常"""
    pass


class PPTGeneratorV2:
    """PPT生成器V2 - 支持丰富的PPT故事线"""

    # 主题色
    PRIMARY_COLOR = RGBColor(46, 80, 144)    # 深蓝色 #2E5090
    ACCENT_COLOR = RGBColor(255, 107, 0)     # 橙色 #FF6B00
    TEXT_COLOR = RGBColor(51, 51, 51)         # 深灰色
    LIGHT_GRAY = RGBColor(245, 245, 245)     # 浅灰色

    def __init__(self, template_path: str = None):
        self.template_path = template_path
        self.presentation = None
        self.template_style = {}  # 存储模板风格

    def create(self, content: Dict[str, Any], language: str = "zh") -> Presentation:
        """根据内容创建PPT"""
        # 提取风格配置
        style_config = content.get("style_config", {})

        if self.template_path and os.path.exists(self.template_path):
            # 加载模板并提取风格，然后创建全新的Presentation（不使用模板文件）
            template_prs = Presentation(self.template_path)
            # 先提取模板风格（在清空之前）
            self._extract_template_style_from_presentation(template_prs)
            # 关闭模板，创建一个全新的空Presentation
            self.presentation = Presentation()
        else:
            self.presentation = Presentation()

        # 设置16:9 (必须在添加幻灯片之前设置)
        self.presentation.slide_width = Inches(13.333)
        self.presentation.slide_height = Inches(7.5)

        # 应用风格配置：优先使用模板提取的风格，然后合并Agent的风格配置
        self._apply_style_config_with_template(style_config)

        slides_data = content.get("slides", [])

        for slide_data in slides_data:
            slide_type = slide_data.get("type", "blank")

            if slide_type == "cover":
                self._add_cover(slide_data)
            elif slide_type == "toc":
                self._add_toc(slide_data)
            elif slide_type == "intro":
                self._add_intro(slide_data)
            elif slide_type == "background":
                self._add_background(slide_data)
            elif slide_type == "overview":
                self._add_overview(slide_data)
            elif slide_type == "dimension_detail":
                self._add_dimension_detail(slide_data)
            elif slide_type == "summary":
                self._add_summary(slide_data)

        return self.presentation

    def _extract_template_style(self):
        """从模板中提取视觉风格 - 兼容版"""
        # 调用新方法，传入self.presentation
        self._extract_template_style_from_presentation(self.presentation)

    def _extract_template_style_from_presentation(self, prs: Presentation):
        """从指定的Presentation中提取视觉风格（包括主题色）"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET

            colors = []
            fonts = []
            theme_colors = {}

            # 1. 从主题文件中提取主题色
            if self.template_path and os.path.exists(self.template_path):
                try:
                    with zipfile.ZipFile(self.template_path, 'r') as z:
                        theme_files = [f for f in z.namelist() if 'theme/theme' in f]
                        if theme_files:
                            xml = z.read(theme_files[0]).decode('utf-8')
                            root = ET.fromstring(xml)
                            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                            color_scheme = root.find('.//a:clrScheme', ns)
                            if color_scheme:
                                for child in color_scheme:
                                    name = child.tag.replace('{http://schemas.openxmlformats.org/drawingml/2006/main}', '')
                                    srgb = child.find('.//a:srgbClr', ns)
                                    if srgb is not None:
                                        val = srgb.get('val', '')
                                        theme_colors[name] = '#' + val
                except Exception as e:
                    pass

            # 2. 从幻灯片文本中提取颜色和字体
            slides = prs.slides._sldIdLst
            for slide_idx in range(min(5, len(slides))):
                try:
                    slide = prs.slides[slide_idx]
                    for shape in slide.shapes:
                        if not hasattr(shape, 'text_frame'):
                            continue
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if not run.font:
                                    continue
                                # 提取颜色
                                if run.font.color:
                                    try:
                                        if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                            c = run.font.color.rgb
                                            colors.append(RGBColor(c.red, c.green, c.blue))
                                    except:
                                        pass
                                # 提取字体
                                if run.font.name:
                                    fonts.append(run.font.name)
                except:
                    pass

            # 3. 选择主色：优先使用主题色dk1，其次从幻灯片提取
            if 'dk1' in theme_colors:
                hex_color = theme_colors['dk1']
                r = int(hex_color[1:3], 16)
                g = int(hex_color[3:5], 16)
                b = int(hex_color[5:7], 16)
                self.template_style['primary_color'] = RGBColor(r, g, b)
            elif colors:
                for c in colors:
                    if not (c.red > 240 and c.green > 240 and c.blue > 240):
                        if not (c.red < 30 and c.green < 30 and c.blue < 30):
                            self.template_style['primary_color'] = c
                            break

            # 4. 设置字体
            if fonts:
                self.template_style['font_name'] = fonts[0]

            # 5. 保存主题色供后续使用
            if theme_colors:
                self.template_style['theme_colors'] = theme_colors

        except Exception as e:
            pass  # 静默失败，使用默认风格

    def _apply_style_config_with_template(self, agent_style_config: Dict):
        """应用风格配置：合并模板风格和Agent风格配置
        优先级：Agent配置 > 模板提取的颜色"""
        # 1. 从模板风格构建配置
        template_style = self.template_style
        combined_config = {"配色": {}, "字体": {}, "布局": {}, "表格样式": {}, "形状样式": {}, "图表样式": {}, "页眉页脚": {}}

        # 从模板提取的主题色（作为默认）
        if 'theme_colors' in template_style:
            combined_config["配色"]["主题色"] = template_style['theme_colors']
            tc = template_style['theme_colors']
            combined_config["配色"]["主色"] = tc.get('dk1', '#2E5090')
            combined_config["配色"]["强调色"] = tc.get('accent1', '#FF6B00')
            combined_config["配色"]["辅色"] = tc.get('lt1', '#F5F5F5')
            combined_config["配色"]["文字色"] = '#333333'

        # 从Agent配置覆盖（Agent优先级最高）
        if agent_style_config:
            agent_colors = agent_style_config.get("配色", {})
            for key in ["主色", "辅色", "强调色", "文字色"]:
                if agent_colors.get(key):
                    combined_config["配色"][key] = agent_colors[key]

            agent_fonts = agent_style_config.get("字体", {})
            if agent_fonts:
                combined_config["字体"] = agent_fonts

            for key in ["布局", "表格样式", "形状样式", "图表样式", "页眉页脚"]:
                if key in agent_style_config:
                    combined_config[key] = agent_style_config[key]

        # 应用合并后的配置
        self._apply_style_config(combined_config)

    def _apply_style_config(self, style_config: Dict):
        """应用完整的风格配置"""
        try:
            # 解析配色
            color_config = style_config.get("配色", {})

            # 解析主色
            primary_hex = color_config.get("主色", "")
            if primary_hex and primary_hex.startswith("#"):
                r = int(primary_hex[1:3], 16)
                g = int(primary_hex[3:5], 16)
                b = int(primary_hex[5:7], 16)
                self.template_style['primary_color'] = RGBColor(r, g, b)

            # 解析辅色
            secondary_hex = color_config.get("辅色", "")
            if secondary_hex and secondary_hex.startswith("#"):
                r = int(secondary_hex[1:3], 16)
                g = int(secondary_hex[3:5], 16)
                b = int(secondary_hex[5:7], 16)
                self.template_style['secondary_color'] = RGBColor(r, g, b)

            # 解析强调色
            accent_hex = color_config.get("强调色", "")
            if accent_hex and accent_hex.startswith("#"):
                r = int(accent_hex[1:3], 16)
                g = int(accent_hex[3:5], 16)
                b = int(accent_hex[5:7], 16)
                self.template_style['accent_color'] = RGBColor(r, g, b)

            # 解析文字色
            text_hex = color_config.get("文字色", "")
            if text_hex and text_hex.startswith("#"):
                r = int(text_hex[1:3], 16)
                g = int(text_hex[3:5], 16)
                b = int(text_hex[5:7], 16)
                self.template_style['text_color'] = RGBColor(r, g, b)

            # 解析字体
            font_config = style_config.get("字体", {})
            self.template_style['font_name'] = font_config.get("主标题", "").split()[0] if font_config.get("主标题") else "微软雅黑"
            self.template_style['title_font'] = font_config.get("主标题", "微软雅黑 Bold 32pt")
            self.template_style['subtitle_font'] = font_config.get("副标题", "微软雅黑 Bold 24pt")
            self.template_style['body_font'] = font_config.get("正文", "微软雅黑 18pt")

            # 解析表格样式
            table_style = style_config.get("表格样式", {})
            if table_style:
                self.template_style['table_title_bg'] = table_style.get("title_bg", "#2E5090")
                self.template_style['table_border'] = table_style.get("border", True)
                self.template_style['table_border_color'] = table_style.get("border_color", "#CCCCCC")

            # 解析形状样式
            shape_style = style_config.get("形状样式", {})
            if shape_style:
                self.template_style['shape_fill_color'] = shape_style.get("fill_color", "#2E5090")
                self.template_style['shape_fill_type'] = shape_style.get("fill_type", "solid")
                self.template_style['shape_line_color'] = shape_style.get("line_color", "#CCCCCC")

            # 解析图表样式
            chart_style = style_config.get("图表样式", {})
            if chart_style:
                self.template_style['chart_type'] = chart_style.get("type", "COLUMN_CLUSTERED")
                self.template_style['chart_has_legend'] = chart_style.get("has_legend", True)
                self.template_style['chart_colors'] = chart_style.get("colors", ["#2E5090", "#FF6B00", "#F5F5F5"])

            # 解析页眉页脚配置
            footer_config = style_config.get("页眉页脚", {})
            if footer_config:
                self.template_style['show_page_num'] = footer_config.get("show_page_num", True)
                self.template_style['show_footer'] = footer_config.get("show_footer", False)

        except Exception as e:
            pass  # 静默失败

    def _get_primary_color(self) -> RGBColor:
        """获取主色，优先使用模板风格"""
        # 优先从theme_colors获取主色
        if 'theme_colors' in self.template_style:
            tc = self.template_style['theme_colors']
            if 'dk1' in tc:
                return self._hex_to_rgb(tc['dk1'])
        if 'primary_color' in self.template_style:
            return self.template_style['primary_color']
        return RGBColor(46, 80, 144)  # 默认深蓝色

    def _get_secondary_color(self) -> RGBColor:
        """获取辅色"""
        if 'theme_colors' in self.template_style:
            tc = self.template_style['theme_colors']
            if 'lt1' in tc:
                return self._hex_to_rgb(tc['lt1'])
        if 'secondary_color' in self.template_style:
            return self.template_style['secondary_color']
        return RGBColor(245, 245, 245)

    def _get_accent_color(self) -> RGBColor:
        """获取强调色"""
        if 'theme_colors' in self.template_style:
            tc = self.template_style['theme_colors']
            if 'accent1' in tc:
                return self._hex_to_rgb(tc['accent1'])
        if 'accent_color' in self.template_style:
            return self.template_style['accent_color']
        return RGBColor(255, 107, 0)

    def _get_text_color(self) -> RGBColor:
        """获取文字颜色，优先使用模板风格"""
        if 'text_color' in self.template_style:
            return self.template_style['text_color']
        return RGBColor(51, 51, 51)  # 默认深灰

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """将十六进制颜色转换为RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    def _add_cover(self, data: Dict):
        """封面页 - 居中布局，背景色块"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

        # 添加背景色块
        bg_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg_box.fill.solid()
        bg_box.fill.fore_color.rgb = self._get_primary_color()
        bg_box.line.fill.background = True

        # 主标题 - 白色
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        tf.text = data.get("title", "云核心网网络架构评估报告")
        tf.paragraphs[0].font.size = Pt(48)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        tf.text = data.get("subtitle", "运营商客户汇报材料")
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 日期 - 底部
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = date_box.text_frame
        tf.text = data.get("date", "2026年4月")
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_toc(self, data: Dict):
        """目录页 - 左侧标题，右侧内容"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

        # 左侧装饰条
        decor = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
        decor.fill.solid()
        decor.fill.fore_color.rgb = self._get_primary_color()
        decor.line.fill.background = True

        # 标题 - 左侧
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(4), Inches(0.8))
        tf = title_box.text_frame
        tf.text = data.get("title", "目录")
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self._get_primary_color()

        # 目录项 - 两列布局
        items = data.get("items", [])
        left_items = items[:6]
        right_items = items[6:]

        # 左侧列
        toc_box1 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(5.5), Inches(5))
        tf = toc_box1.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{i+1}. {item}"
            p.font.size = Pt(16)
            p.space_before = Pt(12)
            if not item.startswith("   "):
                p.font.bold = True

        # 右侧列
        toc_box2 = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.5), Inches(5))
        tf = toc_box2.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{i+len(left_items)+1}. {item}"
            p.font.size = Pt(16)
            p.space_before = Pt(12)
            if not item.startswith("   "):
                p.font.bold = True

    def _add_intro(self, data: Dict):
        """评估方案介绍页 - 左右分栏布局"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self._add_title(slide, data.get("title", "评估方案介绍"))

        content = data.get("content", {})

        # 左侧：评估维度
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(2))
        tf = left_box.text_frame
        p = tf.paragraphs[0]
        p.text = "评估维度"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self._get_primary_color()

        dims = content.get("评估维度", [])
        for dim in dims:
            p = tf.add_paragraph()
            p.text = f"• {dim}"
            p.font.size = Pt(14)
            p.space_before = Pt(6)

        # 右侧：评估产品域
        right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(6), Inches(2))
        tf = right_box.text_frame
        p = tf.paragraphs[0]
        p.text = "评估产品域"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self._get_primary_color()

        domains = content.get("评估产品域", [])
        for domain in domains:
            p = tf.add_paragraph()
            p.text = f"• {domain}"
            p.font.size = Pt(14)
            p.space_before = Pt(6)

        # 底部：评估方法和标准
        bottom_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(12), Inches(2.5))
        tf = bottom_box.text_frame

        p = tf.paragraphs[0]
        p.text = "评估方法"
        p.font.size = Pt(16)
        p.font.bold = True

        method = content.get("评估方法", "")
        p = tf.add_paragraph()
        p.text = method
        p.font.size = Pt(14)
        p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = "评估标准"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(16)

        standard = content.get("评估标准", "")
        p = tf.add_paragraph()
        p.text = standard
        p.font.size = Pt(14)

    def _add_background(self, data: Dict):
        """业务背景页 - 上方统计卡片，下方表格"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self._add_title(slide, data.get("title", "业务背景"))

        content = data.get("content", {})

        # 上方：两个统计卡片并排
        # 卡片1：网元数
        card1 = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(5.5), Inches(1.5))
        card1.fill.solid()
        card1.fill.fore_color.rgb = self._get_primary_color()
        card1.line.fill.background = True

        tf1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5)).text_frame
        tf1.paragraphs[0].text = "总体网元数"
        tf1.paragraphs[0].font.size = Pt(14)
        tf1.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        tf1.paragraphs[0].alignment = PP_ALIGN.CENTER

        tf2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.6)).text_frame
        tf2.paragraphs[0].text = str(content.get('总体网元数', 0))
        tf2.paragraphs[0].font.size = Pt(36)
        tf2.paragraphs[0].font.bold = True
        tf2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 卡片2：容灾组数
        card2 = slide.shapes.add_shape(1, Inches(7), Inches(1.5), Inches(5.5), Inches(1.5))
        card2.fill.solid()
        card2.fill.fore_color.rgb = self._get_text_color()
        card2.line.fill.background = True

        tf3 = slide.shapes.add_textbox(Inches(7), Inches(1.7), Inches(5.5), Inches(0.5)).text_frame
        tf3.paragraphs[0].text = "容灾组数"
        tf3.paragraphs[0].font.size = Pt(14)
        tf3.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        tf3.paragraphs[0].alignment = PP_ALIGN.CENTER

        tf4 = slide.shapes.add_textbox(Inches(7), Inches(2.3), Inches(5.5), Inches(0.6)).text_frame
        tf4.paragraphs[0].text = str(content.get('容灾组数', 0))
        tf4.paragraphs[0].font.size = Pt(36)
        tf4.paragraphs[0].font.bold = True
        tf4.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf4.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 下方：产品域分布表格
        domains = content.get("产品域分布", [])
        if domains:
            self._add_table(slide, 3.5, domains, ["产品域", "网元数", "容灾组数"])

    def _add_overview(self, data: Dict):
        """整体评估概览页 - 顶部结论，表格居中"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

        # 标题带背景色
        title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(13.333), Inches(1))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = self._get_primary_color()
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "整体评估概览")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        content = data.get("content", {})

        # 整体结论
        overall = content.get("overall", "")
        if overall:
            box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.8))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = overall
            p.font.size = Pt(16)
            p.font.bold = True

        # 评估结果表格
        highlights = content.get("highlights", [])
        if highlights:
            self._add_overview_table(slide, 2.2, highlights)

    def _add_overview_table(self, slide, top: float, highlights: List[Dict]):
        """添加评估概览表格"""
        rows = len(highlights) + 1  # +1 header
        cols = 4

        table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(top), Inches(11.5), Inches(0.5 * rows)).table

        # Header
        headers = ["评估维度", "通过率", "风险等级", "状态"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._get_primary_color()
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True

        # Data
        for i, h in enumerate(highlights):
            table.cell(i+1, 0).text = h.get("维度", "")
            table.cell(i+1, 1).text = h.get("通过率", "")
            table.cell(i+1, 2).text = h.get("风险", "")
            table.cell(i+1, 3).text = h.get("状态", "")

            # 风险等级着色
            risk = h.get("风险", "")
            if risk == "高":
                table.cell(i+1, 2).fill.solid()
                table.cell(i+1, 2).fill.fore_color.rgb = RGBColor(255, 200, 200)

    def _add_dimension_detail(self, data: Dict):
        """评估维度详情页 - 顶部色块+三栏布局"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

        # 顶部色带
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = self._get_primary_color()
        header.line.fill.background = True

        # 标题在色带中
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # 获取数据
        content = data.get("content", {}) if isinstance(data.get("content"), dict) else {}
        pass_rate = data.get("pass_rate") or content.get("pass_rate") or ""
        risk = data.get("risk") or content.get("risk") or ""

        # 通过率 - 左侧
        box1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(3.5), Inches(1))
        tf = box1.text_frame
        p = tf.paragraphs[0]
        p.text = f"通过率"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(120, 120, 120)
        p2 = tf.add_paragraph()
        p2.text = pass_rate
        p2.font.size = Pt(32)
        p2.font.bold = True

        # 风险等级 - 中间
        box2 = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(3.5), Inches(1))
        tf = box2.text_frame
        p = tf.paragraphs[0]
        p.text = f"风险等级"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(120, 120, 120)
        p2 = tf.add_paragraph()
        p2.text = risk
        p2.font.size = Pt(32)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(200, 0, 0) if risk == "高" else self.TEXT_COLOR

        # 评估发现 - 左下
        findings = data.get("findings") or content.get("findings", [])
        if findings:
            box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(4), Inches(2))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = "评估发现"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._get_primary_color()

            for f in findings[:3]:
                p = tf.add_paragraph()
                p.text = f"• {f}"
                p.font.size = Pt(11)
                p.space_before = Pt(4)

        # 主要风险 - 中下
        risks = data.get("risks") or content.get("risks", [])
        if risks:
            box = slide.shapes.add_textbox(Inches(5), Inches(2.8), Inches(4), Inches(2))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = "主要风险"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._get_primary_color()

            for r in risks:
                p = tf.add_paragraph()
                p.text = f"⚠ {r}"
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(200, 0, 0)
                p.space_before = Pt(4)

        # 优化建议 - 底部
        suggestions = data.get("suggestions") or content.get("suggestions", [])
        if suggestions:
            box = slide.shapes.add_textbox(Inches(9.2), Inches(2.8), Inches(3.8), Inches(2))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = "优化建议"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._get_primary_color()

            for s in suggestions:
                p = tf.add_paragraph()
                p.text = f"→ {s}"
                p.font.size = Pt(11)
                p.space_before = Pt(4)

    def _add_summary(self, data: Dict):
        """总结页 - 顶部色块+表格+底部"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

        # 顶部色带
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1))
        header.fill.solid()
        header.fill.fore_color.rgb = self._get_primary_color()
        header.line.fill.background = True

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "总结与优化建议")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        content = data.get("content", {})

        # 整体结论
        overall = content.get("overall", "")
        if overall:
            box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.8))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = overall
            p.font.size = Pt(16)
            p.font.bold = True

        # 优先级建议表格
        priorities = content.get("priorities", [])
        if priorities:
            self._add_priority_table(slide, 2.2, priorities)

        # 下一步
        next_steps = content.get("next_steps", "")
        if next_steps:
            box = slide.shapes.add_textbox(Inches(0.8), Inches(6), Inches(11.5), Inches(1))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"下一步: {next_steps}"
            p.font.size = Pt(14)
            p.font.italic = True

    def _add_priority_table(self, slide, top: float, priorities: List[Dict]):
        """添加优先级表格"""
        rows = len(priorities) + 1
        cols = 3

        table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(top), Inches(11.5), Inches(0.4 * rows)).table

        # Header
        headers = ["优先级", "问题", "建议"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._get_primary_color()
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True

        # Data
        for i, p in enumerate(priorities):
            table.cell(i+1, 0).text = p.get("优先级", "")
            table.cell(i+1, 1).text = p.get("问题", "")
            table.cell(i+1, 2).text = p.get("建议", "")

            # 优先级着色
            level = p.get("优先级", "")
            if level == "高":
                table.cell(i+1, 0).fill.solid()
                table.cell(i+1, 0).fill.fore_color.rgb = RGBColor(255, 200, 200)

    def _add_table(self, slide, top: float, data: List[Dict], headers: List[str]):
        """添加通用表格"""
        rows = len(data) + 1
        cols = len(headers)

        table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(top), Inches(11.5), Inches(0.4 * rows)).table

        # Header
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._get_primary_color()
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Data
        for i, row in enumerate(data):
            for j, key in enumerate(headers):
                key_map = {headers[0]: "name", headers[1]: "ne_count", headers[2]: "dr_count"}
                table.cell(i+1, j).text = str(row.get(key_map.get(key, key), ""))

    def _add_title(self, slide, title: str):
        """添加页面标题"""
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self._get_primary_color()

    def _add_chart(self, slide, chart_type: str, data: List, title: str, left: float, top: float, width: float, height: float):
        """
        添加图表（柱状图/折线图/饼图）

        Args:
            slide: 幻灯片对象
            chart_type: 图表类型 (bar/line/pie)
            data: 数据 [categories, values]
            title: 图表标题
            left/top/width/height: 位置和尺寸（英寸）
        """
        try:
            # 映射图表类型
            chart_type_map = {
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
            }
            xl_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

            chart_data = ChartData()
            chart_data.categories = data[0] if len(data) > 0 else []
            chart_data.add_series(title, data[1] if len(data) > 1 else [])

            xlsx_chart = slide.shapes.add_chart(
                xl_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
            ).chart

            return xlsx_chart
        except Exception as e:
            self._safe_print(f"[Chart] Failed to add chart: {e}")
            return None

    def _calc_optimal_font_size(self, text: str, max_width: float = 12) -> float:
        """
        根据文本长度计算最佳字体大小

        Args:
            text: 文本内容
            max_width: 最大宽度（英寸）

        Returns:
            最佳字体大小
        """
        base_size = 24
        if len(text) < 50:
            return base_size
        return max(12, base_size - (len(text) - 50) / 50 * 4)

    def _get_safe_font(self, preferred: str) -> str:
        """获取安全的字体名称（跨平台fallback）"""
        import sys

        font_fallback = {
            "微软雅黑": ["Microsoft YaHei", "SimHei", "Arial"],
            "黑体": ["SimHei", "Microsoft YaHei", "Arial"],
        }

        if sys.platform == "win32":
            return preferred
        return font_fallback.get(preferred, ["Arial"])[0]

    def _safe_print(self, msg: str, limit: int = 200):
        """Safe print to avoid encoding errors"""
        try:
            print(msg[:limit] + "..." if len(msg) > limit else msg)
        except UnicodeEncodeError:
            print("[Output - chars truncated]")

    def save(self, output_path: str) -> str:
        """保存PPT"""
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        try:
            self.presentation.save(output_path)
        except Exception as e:
            raise PPTGeneratorError(f"Failed to save PPT: {e}") from e

        # 修复16:9幻灯片类型属性（python-pptx默认设置为screen4x3）
        try:
            self._fix_slide_size_type(output_path)
        except Exception as e:
            raise PPTGeneratorError(f"Failed to fix slide size: {e}") from e

        return output_path

    def _fix_slide_size_type(self, pptx_path: str):
        """修复PPT的幻灯片尺寸类型为16:9"""
        import zipfile
        import shutil
        try:
            with zipfile.ZipFile(pptx_path, 'r') as z_in:
                with zipfile.ZipFile(pptx_path + '.tmp', 'w') as z_out:
                    for item in z_in.infolist():
                        data = z_in.read(item.filename)
                        if item.filename == 'ppt/presentation.xml':
                            # 修复type属性：screen4x3 -> screen16x9
                            data = data.replace(b'type="screen4x3"', b'type="screen16x9"')
                        z_out.writestr(item, data)
            shutil.move(pptx_path + '.tmp', pptx_path)
        except Exception as e:
            # 清理临时文件
            if os.path.exists(pptx_path + '.tmp'):
                shutil.remove(pptx_path + '.tmp')
            raise PPTGeneratorError(f"Failed to fix slide size: {e}") from e


def generate_ppt(content: Dict[str, Any], output_path: str, template_path: str = None, language: str = "zh") -> str:
    """生成PPT"""
    generator = PPTGeneratorV2(template_path)
    generator.create(content, language)
    return generator.save(output_path)


if __name__ == "__main__":
    # Test
    test_content = {
        "slides": [
            {"type": "cover", "title": "Test Report", "subtitle": "Test", "date": "2026-04"},
            {"type": "toc", "title": "TOC", "items": ["1. Intro", "2. Overview"]},
            {"type": "overview", "title": "Overview", "content": {
                "highlights": [
                    {"维度": "MT0", "通过率": "80%", "风险": "中", "状态": "良好"}
                ],
                "overall": "整体良好"
            }},
            {"type": "summary", "title": "Summary", "content": {
                "overall": "测试完成",
                "priorities": [{"优先级": "高", "问题": "test", "建议": "fix"}]
            }}
        ]
    }
    result = generate_ppt(test_content, "test.pptx")
    print(f"Generated: {result}")