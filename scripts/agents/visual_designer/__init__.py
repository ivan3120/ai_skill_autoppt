# -*- coding: utf-8 -*-
"""
VisualDesigner Agent - 视觉设计师

独立解耦的Agent模块，负责分析PPT模板风格并提取视觉特征。
"""

import os
from typing import Dict
from scripts.agents.base import BaseAgent, AgentResult


class VisualDesignerAgent(BaseAgent):
    """
    视觉设计师Agent
    任务：分析用户提供的PPT模板风格，提取视觉特征，将风格应用到生成的PPT中
    """

    def __init__(self):
        super().__init__("VisualDesigner", "视觉设计师", "visual_designer")
        # 设置默认prompt模板
        if not self.prompt_template:
            self.prompt_template = """你是一个PPT视觉设计专家。请分析以下模板风格信息，生成完整的风格配置。

模板主题色信息：
- dk1 (深色1): {dk1}
- lt1 (浅色1): {lt1}
- accent1 (强调色1): {accent1}
- accent2 (强调色2): {accent2}
- accent3 (强调色3): {accent3}
- accent4 (强调色4): {accent4}
- accent5 (强调色5): {accent5}
- accent6 (强调色6): {accent6}

请返回JSON格式的风格配置：
```json
{
    "配色": {
        "主色": "主色调 hex",
        "辅色": "辅色 hex",
        "强调色": "强调色 hex",
        "文字色": "文字色 hex"
    },
    "字体": {
        "主标题": "字体名 字号",
        "副标题": "字体名 字号",
        "正文": "字体名 字号"
    },
    "表格样式": {...},
    "图表样式": {...}
}
```"""

    def execute(self, input_data: Dict) -> AgentResult:
        """
        学习并应用模板风格 - 混合方式（代码提取+LLM分析）
        """
        template_path = input_data.get("template_path")

        if template_path and os.path.exists(template_path):
            # Step 1: 用代码提取基础信息
            basic_info = self._extract_basic_style(template_path)
            self._safe_print(f"[VisualDesigner] Basic info extracted: {len(basic_info.get('colors', []))} colors, {len(basic_info.get('fonts', []))} fonts")

            # Step 2: 用LLM进行深度分析
            if self._has_llm():
                style_config = self._llm_analyze_template(template_path, basic_info)
                self._safe_print(f"[VisualDesigner] LLM analysis completed")
            else:
                style_config = self._convert_basic_to_config(basic_info)
        else:
            style_config = {
                "配色": {"主色": "#2E5090", "辅色": "#F5F5F5", "强调色": "#FF6B00", "文字色": "#333333"},
                "字体": {"主标题": "微软雅黑 Bold 32pt", "副标题": "微软雅黑 Bold 24pt", "正文": "微软雅黑 Regular 18pt"},
                "布局": {"标题位置": "顶部居中", "内容边距": "0.5英寸"},
                "页眉页脚": {"显示页码": True, "页码位置": "底部居中"}
            }

        return AgentResult(
            status="success",
            data={"style_config": style_config},
            message="已学习模板风格"
        )

    def _extract_basic_style(self, template_path: str) -> Dict:
        """用代码提取PPT模板的完整样式信息"""
        try:
            from pptx import Presentation
            import zipfile
            import re
            import xml.etree.ElementTree as ET

            prs = Presentation(template_path)
            result = {
                "slide_count": len(prs.slides),
                "colors": [], "fonts": [], "font_sizes": [],
                "bg_colors": [], "table_styles": [], "theme_colors": {},
                "shape_styles": [], "chart_styles": [], "footer_config": {}
            }

            # 提取主题颜色
            try:
                with zipfile.ZipFile(template_path, 'r') as z:
                    theme_files = [f for f in z.namelist() if 'theme/theme' in f]
                    if theme_files:
                        xml = z.read(theme_files[0]).decode('utf-8')
                        root = ET.fromstring(xml)
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        color_scheme = root.find('.//a:clrScheme', ns)
                        if color_scheme:
                            for child in color_scheme:
                                name = child.tag.replace('{http://schemas.openxmlformats.org/drawingml/2006/main}', '')
                                srgb = child.find('.//a:srgbClr', ns)
                                if srgb is not None:
                                    result["theme_colors"][name] = f"#{srgb.get('val', '')}"
            except Exception as e:
                pass

            slides_list = prs.slides._sldIdLst
            for slide_idx in range(min(10, len(slides_list))):
                try:
                    slide = prs.slides[slide_idx]

                    # 提取背景色
                    try:
                        if slide.background and slide.background.fill:
                            fc = slide.background.fill.fore_color
                            if hasattr(fc, 'rgb'):
                                result["bg_colors"].append(f"#{fc.rgb.red:02X}{fc.rgb.green:02X}{fc.rgb.blue:02X}")
                    except:
                        pass

                    # 提取页眉页脚信息
                    try:
                        if slide.placeholders:
                            for ph in slide.placeholders:
                                if ph.placeholder_format.type == 8:  # 页脚
                                    result["footer_config"]["has_footer"] = True
                                elif ph.placeholder_format.type == 9:  # 页码
                                    result["footer_config"]["has_page_num"] = True
                    except:
                        pass

                    # 遍历所有形状
                    for shape in slide.shapes:
                        # 表格样式提取
                        if hasattr(shape, "has_table") and shape.has_table:
                            table = shape.table
                            if table.rows and table.columns:
                                first_row = table.rows[0]
                                cells_bg = []
                                cell_fonts = []
                                for cell in first_row.cells:
                                    try:
                                        if cell.fill.fore_color and hasattr(cell.fill.fore_color, 'rgb'):
                                            c = cell.fill.fore_color.rgb
                                            cells_bg.append(f"#{c.red:02X}{c.green:02X}{c.blue:02X}")
                                        if cell.text_frame and cell.text_frame.paragraphs:
                                            for run in cell.text_frame.paragraphs[0].runs:
                                                if run.font and run.font.name:
                                                    cell_fonts.append(run.font.name)
                                    except:
                                        pass
                                if cells_bg:
                                    result["table_styles"].append({
                                        "title_bg": cells_bg[0],
                                        "font": cell_fonts[0] if cell_fonts else ""
                                    })

                        # 形状样式提取
                        if hasattr(shape, "fill") and shape.fill:
                            try:
                                shape_info = {}
                                if shape.fill.fore_color and hasattr(shape.fill.fore_color, 'rgb'):
                                    c = shape.fill.fore_color.rgb
                                    shape_info["fill_color"] = f"#{c.red:02X}{c.green:02X}{c.blue:02X}"
                                shape_info["fill_type"] = str(shape.fill.fill_type) if shape.fill.fill_type else "solid"
                                if hasattr(shape, "line") and shape.line and shape.line.fill:
                                    if shape.line.fill.fore_color and hasattr(shape.line.fill.fore_color, 'rgb'):
                                        lc = shape.line.fill.fore_color.rgb
                                        shape_info["line_color"] = f"#{lc.red:02X}{lc.green:02X}{lc.blue:02X}"
                                if shape_info:
                                    result["shape_styles"].append(shape_info)
                            except:
                                pass

                        # 图表样式提取
                        if hasattr(shape, "has_chart") and shape.has_chart:
                            try:
                                chart = shape.chart
                                chart_info = {
                                    "type": str(chart.chart_type),
                                    "has_legend": chart.has_legend if hasattr(chart, 'has_legend') else False
                                }
                                result["chart_styles"].append(chart_info)
                            except:
                                pass

                        # 文本样式提取
                        if not hasattr(shape, "text_frame"):
                            continue
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if not run.font:
                                    continue
                                try:
                                    if run.font.color and hasattr(run.font.color, 'rgb'):
                                        c = run.font.color.rgb
                                        result["colors"].append(f"#{c.red:02X}{c.green:02X}{c.blue:02X}")
                                except:
                                    pass
                                if run.font.name:
                                    result["fonts"].append(run.font.name)
                                if run.font.size:
                                    result["font_sizes"].append(int(run.font.size.pt))

                except Exception as e:
                    self._safe_print(f"[VisualDesigner] Slide {slide_idx} extraction error: {e}")
                    continue

            result["colors"] = list(set(result["colors"]))
            result["fonts"] = list(set(result["fonts"]))
            result["font_sizes"] = sorted(list(set(result["font_sizes"])))
            result["bg_colors"] = list(set(result["bg_colors"]))

            return result
        except Exception as e:
            self._safe_print(f"[VisualDesigner] Basic extraction failed: {e}")
            return {}

    def _llm_analyze_template(self, template_path: str, basic_info: Dict) -> Dict:
        """用LLM深度分析模板风格"""
        theme_colors = basic_info.get("theme_colors", {})

        user_message = self._load_llm_prompt(
            self.llm_prompt_file,
            dk1=theme_colors.get('dk1', 'N/A'),
            lt1=theme_colors.get('lt1', 'N/A'),
            accent1=theme_colors.get('accent1', 'N/A'),
            accent2=theme_colors.get('accent2', 'N/A'),
            accent3=theme_colors.get('accent3', 'N/A'),
            accent4=theme_colors.get('accent4', 'N/A'),
            accent5=theme_colors.get('accent5', 'N/A'),
            accent6=theme_colors.get('accent6', 'N/A')
        )

        try:
            llm_result = self.invoke_llm(user_message)
            return self._parse_llm_style_result(llm_result)
        except Exception as e:
            self._safe_print(f"[VisualDesigner] LLM analysis failed: {e}")
            return self._convert_basic_to_config(basic_info)

    def _parse_llm_style_result(self, llm_result: str) -> Dict:
        """解析LLM返回的风格配置"""
        import json
        try:
            if "```json" in llm_result:
                start = llm_result.find("```json") + 7
                end = llm_result.find("```", start)
                return json.loads(llm_result[start:end].strip())
            elif "{" in llm_result:
                start = llm_result.find("{")
                end = llm_result.rfind("}") + 1
                return json.loads(llm_result[start:end])
        except Exception as e:
            self._safe_print(f"[VisualDesigner] JSON parse failed: {e}")

        return {
            "配色": {"主色": "#2E5090", "辅色": "#F5F5F5", "强调色": "#FF6B00", "文字色": "#333333"},
            "字体": {"主标题": "微软雅黑 Bold 32pt", "副标题": "微软雅黑 Bold 24pt", "正文": "微软雅黑 18pt"}
        }

    def _convert_basic_to_config(self, basic_info: Dict) -> Dict:
        """将基础信息转换为完整的风格配置"""
        colors = basic_info.get("colors", [])
        fonts = basic_info.get("fonts", [])
        sizes = basic_info.get("font_sizes", [])
        table_styles = basic_info.get("table_styles", [])
        shape_styles = basic_info.get("shape_styles", [])
        chart_styles = basic_info.get("chart_styles", [])
        footer_config = basic_info.get("footer_config", {})
        theme_colors = basic_info.get("theme_colors", {})

        # 主色：优先使用主题色dk1，否则使用最常见的非黑白颜色
        main_color = theme_colors.get('dk1') or "#2E5090"
        for c in colors:
            if c not in ["#FFFFFF", "#000000", "#ffffff", "#000000"]:
                main_color = c
                break

        # 辅色：使用主题色lt1
        secondary_color = theme_colors.get('lt1') or "#F5F5F5"

        # 强调色：使用主题色accent1
        accent_color = theme_colors.get('accent1') or "#FF6B00"

        # 主字体
        main_font = fonts[0] if fonts else "微软雅黑"

        # 标题/正文字体大小
        title_size = 32
        body_size = 18
        if sizes:
            title_size = max([s for s in sizes if s <= 44], default=32)
            body_size = min([s for s in sizes if s >= 12], default=18)

        # 表格样式
        table_style = {}
        if table_styles:
            ts = table_styles[0]
            table_style = {
                "title_bg": ts.get("title_bg", "#2E5090"),
                "title_font": ts.get("font", main_font),
                "border": True,
                "border_color": "#CCCCCC"
            }

        # 形状样式
        shape_style = {}
        if shape_styles:
            ss = shape_styles[0]
            shape_style = {
                "fill_color": ss.get("fill_color", main_color),
                "fill_type": ss.get("fill_type", "solid"),
                "line_color": ss.get("line_color", "#CCCCCC")
            }

        # 图表样式
        chart_style = {}
        if chart_styles:
            cs = chart_styles[0]
            chart_style = {
                "type": cs.get("type", "COLUMN_CLUSTERED"),
                "has_legend": cs.get("has_legend", True),
                "colors": [main_color, accent_color, secondary_color]
            }

        # 页眉页脚配置
        footer_style = {
            "show_page_num": footer_config.get("has_page_num", True),
            "show_footer": footer_config.get("has_footer", False),
            "page_num_position": "底部居中",
            "footer_position": "底部左侧"
        }

        return {
            "配色": {
                "主色": main_color,
                "辅色": secondary_color,
                "强调色": accent_color,
                "文字色": "#333333",
                "主题色": theme_colors
            },
            "字体": {
                "主标题": f"{main_font} Bold {title_size}pt",
                "副标题": f"{main_font} Bold 24pt",
                "正文": f"{main_font} {body_size}pt"
            },
            "布局": {
                "标题位置": "顶部居中",
                "内容边距": "0.5英寸"
            },
            "表格样式": table_style,
            "形状样式": shape_style,
            "图表样式": chart_style,
            "页眉页脚": footer_style
        }
